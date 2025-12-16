import streamlit as st
import os
import time  # Add time import for latency measurement
from rag_chatbot import ask, general_ask, vectordb  # Import ask and general_ask functions and vectordb
from langchain_core.documents import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter
import pytesseract
from PIL import Image
import tempfile
from langchain_community.document_loaders import PyPDFLoader
from pptx import Presentation

# Configure Tesseract path
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

st.set_page_config(page_title="RAG Chatbot", page_icon="📚")

st.title(" Chatbot Pro is ready to go")
st.write("Ask questions based on the knowledge base. You can also upload PDFs, images, or PowerPoint files to add to the knowledge base.")

# Mode selection
mode = st.radio("Select mode:", ["RAG (strictly from knowledge base)", "General (may hallucinate)"], index=0)

# File uploader
uploaded_file = st.file_uploader("Upload PDF, JPG, PNG, PPTX files", type=["pdf", "jpg", "png", "pptx"])

if uploaded_file is not None:
    with st.spinner("Processing uploaded file..."):
        try:
            if uploaded_file.type.startswith("image"):
                # Process image
                img = Image.open(uploaded_file)
                text = pytesseract.image_to_string(img)
                if text.strip():
                    doc = Document(page_content=text, metadata={"source": uploaded_file.name})
                    splitter = RecursiveCharacterTextSplitter(chunk_size=800, chunk_overlap=150)
                    chunks = splitter.split_documents([doc])
                    vectordb.add_documents(chunks)
                    st.success(f"Image '{uploaded_file.name}' processed and added to knowledge base.")
                else:
                    st.warning("No text found in the image.")
            elif uploaded_file.type == "application/pdf":
                # Process PDF
                with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_file:
                    temp_file.write(uploaded_file.getbuffer())
                    temp_path = temp_file.name
                
                loader = PyPDFLoader(temp_path)
                docs = loader.load()
                for doc in docs:
                    doc.metadata["source"] = uploaded_file.name
                os.unlink(temp_path)
                if docs:
                    splitter = RecursiveCharacterTextSplitter(chunk_size=800, chunk_overlap=150)
                    chunks = splitter.split_documents(docs)
                    vectordb.add_documents(chunks)
                    st.success(f"PDF '{uploaded_file.name}' processed and added to knowledge base.")
                else:
                    st.warning("No content found in the PDF.")
            elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                # Process PowerPoint
                with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as temp_file:
                    temp_file.write(uploaded_file.getbuffer())
                    temp_path = temp_file.name
                
                prs = Presentation(temp_path)
                text = ""
                for slide_num, slide in enumerate(prs.slides, start=1):
                    slide_text = f"Slide {slide_num}:\n"
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            slide_text += shape.text + "\n"
                    text += slide_text + "\n"
                os.unlink(temp_path)
                if text.strip():
                    doc = Document(page_content=text, metadata={"source": uploaded_file.name})
                    splitter = RecursiveCharacterTextSplitter(chunk_size=800, chunk_overlap=150)
                    chunks = splitter.split_documents([doc])
                    vectordb.add_documents(chunks)
                    st.success(f"PowerPoint '{uploaded_file.name}' processed and added to knowledge base.")
                else:
                    st.warning("No content found in the PowerPoint file.")
        except Exception as e:
            st.error(f"Error processing file: {e}")

        except Exception as e:
            st.error(f"Error processing file: {e}")

# Initialize session state for history
if "history" not in st.session_state:
    st.session_state.history = []

# Input box
user_input = st.text_input("Ask a question:", "", key="input")

# When user clicks button
if st.button("Ask"):
    if user_input.strip() == "":
        st.warning("Please enter a question.")
    else:
        start_time = time.time()  # Start timing
        
        if mode == "RAG (strictly from knowledge base)":
            answer, citations = ask(user_input, return_data=True)
            # Filter citations to only from data/
            filtered_citations = [(src, pg) for src, pg in citations if isinstance(src, str) and src.startswith("data")]
            # Remove duplicates
            unique_citations = list(dict.fromkeys(filtered_citations))
            # Citation accuracy: percentage of retrieved docs from knowledge base
            total_retrieved = len(citations)
            accurate_citations = len(unique_citations)
            citation_accuracy = (accurate_citations / total_retrieved * 100) if total_retrieved > 0 else 0
        else:
            answer = general_ask(user_input)
            unique_citations = []
            citation_accuracy = None  # No citations in general mode
            total_retrieved = 0
            accurate_citations = 0
        
        end_time = time.time()  # End timing
        latency = end_time - start_time

        # Add to history
        st.session_state.history.append((user_input, answer, unique_citations))
        if len(st.session_state.history) > 5:
            st.session_state.history.pop(0)

        st.subheader("Answer")
        st.write(answer)

        if unique_citations:
            st.subheader("Citations")
            for src, pg in unique_citations:
                if pg == "N/A":
                    st.write(f"- **{src}** (from image or slide)")
                else:
                    st.write(f"- **{src}**, page {pg}")

        # Display metrics
        st.subheader("System Metrics")
        st.write(f"**Latency**: {latency:.2f} seconds")
        if citation_accuracy is not None:
            st.write(f"**Citation Accuracy**: {citation_accuracy:.1f}% ({accurate_citations}/{total_retrieved} citations from knowledge base)")
        else:
            st.write("**Citation Accuracy**: N/A (General mode)")

# Display history
if st.session_state.history:
    st.subheader("Conversation History (last 5)")
    for i, (q, a, cits) in enumerate(st.session_state.history):
        st.write(f"**Q{i+1}:** {q}")
        st.write(f"**A{i+1}:** {a}")
        if cits:
            st.write("Citations:")
            for src, pg in cits:
                st.write(f"  - {src}, page {pg}")
        st.write("---")

# Clear history button
if st.button("Clear Conversation History"):
    st.session_state.history = []
    st.success("Conversation history cleared!")
