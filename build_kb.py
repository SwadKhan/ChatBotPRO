import os
import glob
from langchain_community.document_loaders import PyPDFLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.embeddings.fastembed import FastEmbedEmbeddings
from langchain_community.vectorstores import Chroma
from langchain_core.documents import Document
import pytesseract
from PIL import Image
from pptx import Presentation

# Configure Tesseract path
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

DATA_DIR = "data"
CHROMA_DIR = "chroma_db"

# Load ALL PDFs from data/ folder and split into chunks
def load_all_pdfs():
    pdf_files = glob.glob(f"{DATA_DIR}/*.pdf")
    all_chunks = []

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=800,
        chunk_overlap=150
    )

    for pdf in pdf_files:
        print(f"Loading: {pdf}")

        loader = PyPDFLoader(pdf)
        docs = loader.load()

        chunks = splitter.split_documents(docs)
        all_chunks.extend(chunks)

        print(f"  -> Pages: {len(docs)}, Chunks: {len(chunks)}")

    print(f"\nTotal chunks from ALL PDFs: {len(all_chunks)}")
    return all_chunks


# Load ALL images from data/ folder
def load_images():
    image_files = glob.glob(f"{DATA_DIR}/*.jpg") + glob.glob(f"{DATA_DIR}/*.png") + glob.glob(f"{DATA_DIR}/*.jpeg")
    documents = []

    for img_file in image_files:
        print(f"Loading image: {img_file}")
        try:
            img = Image.open(img_file)
            text = pytesseract.image_to_string(img)
            if text.strip():
                doc = Document(page_content=text, metadata={"source": img_file})
                documents.append(doc)
                print(f"  -> Extracted text from {img_file}")
            else:
                print(f"  -> No text found in {img_file}")
        except Exception as e:
            print(f"  -> Error processing {img_file}: {e}")

    print(f"Total documents from images: {len(documents)}")
    return documents


# Load ALL videos from data/ folder
def load_videos():
    video_files = glob.glob(f"{DATA_DIR}/*.mp4") + glob.glob(f"{DATA_DIR}/*.avi") + glob.glob(f"{DATA_DIR}/*.mov")
    documents = []

    for vid_file in video_files:
        print(f"Loading video: {vid_file}")
        try:
            cap = cv2.VideoCapture(vid_file)
            fps = cap.get(cv2.CAP_PROP_FPS)
            if fps == 0:
                fps = 30  # default
            frame_interval = int(fps * 30)  # every 30 seconds
            frame_count = 0
            frame_num = 0

            while True:
                ret, frame = cap.read()
                if not ret:
                    break
                if frame_count % frame_interval == 0:
                    img = Image.fromarray(cv2.cvtColor(frame, cv2.COLOR_BGR2RGB))
                    text = pytesseract.image_to_string(img)
                    if text.strip():
                        doc = Document(page_content=text, metadata={"source": vid_file, "frame": frame_num})
                        documents.append(doc)
                        print(f"  -> Extracted text from frame {frame_num} of {vid_file}")
                    frame_num += 1
                frame_count += 1

            cap.release()
        except Exception as e:
            print(f"  -> Error processing {vid_file}: {e}")

    print(f"Total documents from videos: {len(documents)}")
    return documents


# Load ALL PowerPoint files from data/ folder
def load_slides():
    pptx_files = glob.glob(f"{DATA_DIR}/*.pptx") + glob.glob(f"{DATA_DIR}/*.ppt")
    documents = []

    for file in pptx_files:
        print(f"Loading PowerPoint: {file}")
        try:
            prs = Presentation(file)
            text = ""
            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_text = f"Slide {slide_num}:\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + "\n"
                text += slide_text + "\n"
            if text.strip():
                doc = Document(page_content=text, metadata={"source": file, "slide_count": len(prs.slides)})
                documents.append(doc)
                print(f"  -> Processed {file}")
        except Exception as e:
            print(f"  -> Error processing {file}: {e}")

    print(f"Total documents from PowerPoint: {len(documents)}")
    return documents


# Load ALL documents (PDFs, images, videos)
def load_all_documents():
    all_docs = []

    # Load PDFs
    pdf_docs = load_all_pdfs()
    all_docs.extend(pdf_docs)

    # Load images
    img_docs = load_images()
    all_docs.extend(img_docs)

    # Load videos
    vid_docs = load_videos()
    all_docs.extend(vid_docs)

    # Load PowerPoint
    slide_docs = load_slides()
    all_docs.extend(slide_docs)

    # Split text documents into chunks
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=800,
        chunk_overlap=150
    )

    # PDFs are already chunks, images and videos are full text, so split them
    non_pdf_docs = [doc for doc in all_docs if not doc.metadata.get("source", "").endswith(".pdf")]
    pdf_chunks = [doc for doc in all_docs if doc.metadata.get("source", "").endswith(".pdf")]

    chunked_docs = splitter.split_documents(non_pdf_docs)
    all_chunks = pdf_chunks + chunked_docs

    print(f"\nTotal chunks from ALL documents: {len(all_chunks)}")
    return all_chunks


# Build the vector database (ChromaDB)
def build_vector_db():
    chunks = load_all_documents()

    embeddings = FastEmbedEmbeddings()

    vectordb = Chroma.from_documents(
        documents=chunks,
        embedding=embeddings,
        persist_directory=CHROMA_DIR
    )

    vectordb.persist()
    print("\nVector DB created successfully!")


if __name__ == "__main__":
    build_vector_db()
